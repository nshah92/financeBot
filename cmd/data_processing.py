import os
import glob
import pickle
import pandas as pd
import pytesseract
import PyPDF2
from pptx import Presentation
from PIL import Image
import vector_search
from py2neo import Graph, Node

graph = Graph("bolt://localhost:7687", auth=("neo4j", "admin123"))

image_folder = "extracted_images"
if not os.path.exists(image_folder):
    os.makedirs(image_folder)

def extract_text_from_pdf(file_path):
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    extracted_texts = []

    for slide_num, slide in enumerate(prs.slides):
        slide_text = ""

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"

        extracted_texts.append(slide_text.strip())

        for shape_index, shape in enumerate(slide.shapes):
            if hasattr(shape, "image"):
                image_bytes = shape.image.blob
                image_path = os.path.join(image_folder, f"slide_{slide_num}_img_{shape_index}.png")

                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)

                image_text = pytesseract.image_to_string(Image.open(image_path))
                extracted_texts.append(image_text.strip())

    return "\n".join(extracted_texts)

def process_csv_for_faiss(file_path):
    df = pd.read_csv(file_path)
    extracted_texts = []

    for _, row in df.iterrows():
        sentence = (
            f"For the tax year {row['Tax Year']}, a {row['Taxpayer Type']} in {row['State']} "
            f"earned ${float(row['Income']):,.2f} from {row['Income Source']}, with deductions of ${float(row['Deductions']):,.2f} "
            f"under {row['Deduction Type']}. The taxable income was ${float(row['Taxable Income']):,.2f} "
            f"at a tax rate of {float(row['Tax Rate'])*100:.2f}%, resulting in ${float(row['Tax Owed']):,.2f} in taxes owed."
        )
        sentence = sentence.replace("\n", " ").replace("\r", " ").strip()
        extracted_texts.append(sentence)

    return extracted_texts

def process_csv_for_neo4j(file_path):
    df = pd.read_csv(file_path)
    for _, row in df.iterrows():
        tax_node = Node(
            "TaxRecord",
            taxpayer_type=row["Taxpayer Type"],
            tax_year=row["Tax Year"],
            transaction_date=row["Transaction Date"],
            income_source=row["Income Source"],
            deduction_type=row["Deduction Type"],
            state=row["State"],
            income=row["Income"],
            deductions=row["Deductions"],
            taxable_income=row["Taxable Income"],
            tax_rate=row["Tax Rate"],
            tax_owed=row["Tax Owed"]
        )
        graph.create(tax_node)


if __name__ == "__main__":
    data_directory = "data"

    # Process PDFs
    print("Processing PDFs...")
    pdf_files = glob.glob(os.path.join(data_directory, "*.pdf"))
    pdf_texts = [extract_text_from_pdf(f) for f in pdf_files]

    # Process PPTX
    print("Processing PPTX files...")
    pptx_files = glob.glob(os.path.join(data_directory, "*.pptx"))
    pptx_texts = [extract_text_from_pptx(f) for f in pptx_files]

    # Process CSVs
    print("Processing CSV files...")
    csv_files = glob.glob(os.path.join(data_directory, "*.csv"))
    csv_texts = []
    for csv_file in csv_files:
        csv_texts.extend(process_csv_for_faiss(csv_file))
        process_csv_for_neo4j(csv_file)

    if pdf_texts:
        print("Indexing PDFs...")
        pdf_embeddings = vector_search.generate_embeddings(pdf_texts)
        pdf_index = vector_search.create_faiss_index(pdf_embeddings)
        vector_search.save_faiss_index(pdf_index, "generated/pdf_index.bin")
        with open("generated/pdf_texts.pkl", "wb") as f:
            pickle.dump(pdf_texts, f)

    if pptx_texts:
        print("Indexing PPTX...")
        pptx_embeddings = vector_search.generate_embeddings(pptx_texts)
        pptx_index = vector_search.create_faiss_index(pptx_embeddings)
        vector_search.save_faiss_index(pptx_index, "generated/pptx_index.bin")
        with open("generated/pptx_texts.pkl", "wb") as f:
            pickle.dump(pptx_texts, f)

    if csv_texts:
        print("Indexing CSVs...")
        csv_embeddings = vector_search.generate_embeddings(csv_texts)
        csv_index = vector_search.create_faiss_index(csv_embeddings)
        vector_search.save_faiss_index(csv_index, "generated/csv_index.bin")
        with open("generated/csv_texts.pkl", "wb") as f:
            pickle.dump(csv_texts, f)

    print("🚀 Data processing completed! Separate FAISS indexes created.")
